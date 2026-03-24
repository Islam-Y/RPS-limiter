#!/usr/bin/env python3
from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from build_predefense_pptx import (
    BODY_COLOR,
    BODY_FONT,
    MUTED_COLOR,
    NAVY,
    TITLE_COLOR,
    TITLE_FONT,
    SlideSpec,
    add_body_box,
    add_caption,
    add_fitted_picture,
    add_table,
    add_text_block,
    clean_text,
    emphasize_box,
    parse_source,
)


SOURCE_DEFAULT = Path("/Users/islam/IdeaProjects/RPS-limiter/services/Презентация_предзащита_5мин.txt")
TEMPLATE_DEFAULT = Path("/Users/islam/Downloads/Shablon_prezentatsii_ITMO_na_russkom_yazyke_1.pptx")
OUTPUT_DEFAULT = Path("/Users/islam/IdeaProjects/RPS-limiter/services/Презентация_предзащита_5мин_ITMO.pptx")

BASE_WIDE_WIDTH = 13.333
BASE_WIDE_HEIGHT = 7.5
TEMPLATE_WIDTH = 10.0
TEMPLATE_HEIGHT = 5.625
SX = TEMPLATE_WIDTH / BASE_WIDE_WIDTH
SY = TEMPLATE_HEIGHT / BASE_WIDE_HEIGHT


def sx(value: float) -> float:
    return value * SX


def sy(value: float) -> float:
    return value * SY


def clear_all_template_slides(prs: Presentation) -> None:
    for index in reversed(range(len(prs.slides))):
        slide_id = prs.slides._sldIdLst[index]
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        del prs.slides._sldIdLst[index]


def remove_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def get_title_placeholder(slide):
    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.type == PP_PLACEHOLDER.TITLE:
            return placeholder
    return None


def get_non_title_placeholders(slide):
    return [
        placeholder
        for placeholder in slide.placeholders
        if placeholder.placeholder_format.type != PP_PLACEHOLDER.TITLE
    ]


def prepare_slide(prs: Presentation, layout_idx: int, title: str, *, remove_non_title: bool = True):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    title_ph = get_title_placeholder(slide)
    if title_ph is not None:
        title_ph.text = clean_text(title)
    if remove_non_title:
        for placeholder in list(get_non_title_placeholders(slide)):
            remove_shape(placeholder)
    return slide


def add_slide_number(slide, number: int, total: int) -> None:
    box = slide.shapes.add_textbox(Inches(9.15), Inches(5.24), Inches(0.55), Inches(0.18))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{number}/{total}"
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = BODY_FONT
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED_COLOR


def add_template_title_slide(prs: Presentation, spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_ph = get_title_placeholder(slide)
    if title_ph is not None:
        title_ph.text = clean_text(spec.title)

    body_ph = next((ph for ph in slide.placeholders if ph.placeholder_format.type != PP_PLACEHOLDER.TITLE), None)
    if body_ph is not None:
        tf = body_ph.text_frame
        tf.clear()
        tf.word_wrap = True
        for idx, line in enumerate(spec.text_lines):
            text = clean_text(line)
            if not text:
                continue
            paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            paragraph.text = text
            paragraph.font.name = BODY_FONT
            paragraph.font.size = Pt(20 if idx == 0 else 16)
            paragraph.font.bold = idx == 0
            paragraph.font.color.rgb = BODY_COLOR
            paragraph.space_after = Pt(3)

    tag = slide.shapes.add_shape(
        5,  # rounded rectangle
        Inches(6.85),
        Inches(4.28),
        Inches(1.75),
        Inches(0.45),
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = NAVY
    tag.line.fill.background()
    tf = tag.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Предзащита, 5 минут"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = BODY_FONT
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    add_slide_number(slide, spec.number, spec.total)


def add_text_image_slide(prs: Presentation, spec: SlideSpec) -> None:
    slide = prepare_slide(prs, 4, spec.title, remove_non_title=True)

    add_body_box(slide, sx(0.55), sy(1.02), sx(5.05), sy(5.7))
    lines = [line for line in spec.text_lines if clean_text(line)]
    if spec.short_line:
        lines.append("")
        lines.append(spec.short_line)
    add_text_block(
        slide,
        lines,
        sx(0.72),
        sy(1.18),
        sx(4.72),
        sy(5.2),
        font_size=16,
        bold_predicate=lambda text: text.endswith(":"),
    )

    if spec.images:
        image = spec.images[0]
        add_fitted_picture(slide, image.path, sx(5.85), sy(1.18), sx(6.92), sy(4.95))
        add_caption(slide, image.caption, sx(5.95), sy(6.15), sx(6.72))

    add_slide_number(slide, spec.number, spec.total)


def add_goal_slide(prs: Presentation, spec: SlideSpec) -> None:
    slide = prepare_slide(prs, 4, spec.title, remove_non_title=True)

    goal_lines = []
    tasks = []
    target = goal_lines
    for line in spec.text_lines:
        text = clean_text(line)
        if not text:
            continue
        if text == "Задачи:":
            target = tasks
            continue
        if text == "Цель работы:":
            continue
        target.append(text)

    goal_text = " ".join(goal_lines[:1])
    emphasize_box(slide, f"Цель работы: {goal_text}", sx(0.68), sy(1.15), sx(11.95), sy(0.9))

    add_body_box(slide, sx(0.68), sy(2.28), sx(5.85), sy(3.95))
    add_body_box(slide, sx(6.8), sy(2.28), sx(5.85), sy(3.95))
    add_text_block(slide, tasks[:4], sx(0.9), sy(2.52), sx(5.35), sy(3.4), font_size=13)
    add_text_block(slide, tasks[4:], sx(7.02), sy(2.52), sx(5.35), sy(3.4), font_size=13)
    emphasize_box(
        slide,
        "Все задачи связаны с практической проверкой решения на реальном проектном стенде.",
        sx(2.1),
        sy(6.32),
        sx(9.0),
        sy(0.44),
    )

    add_slide_number(slide, spec.number, spec.total)


def add_method_slide(prs: Presentation, spec: SlideSpec) -> None:
    slide = prepare_slide(prs, 4, spec.title, remove_non_title=True)

    add_body_box(slide, sx(0.55), sy(1.05), sx(5.3), sy(5.7))
    add_text_block(
        slide,
        [line for line in spec.text_lines if clean_text(line)],
        sx(0.72),
        sy(1.18),
        sx(4.95),
        sy(5.1),
        font_size=15,
        bold_predicate=lambda text: text.endswith(":"),
    )
    emphasize_box(slide, "Норма → атака → восстановление", sx(0.9), sy(6.05), sx(4.6), sy(0.4))

    if spec.images:
        image = spec.images[0]
        add_fitted_picture(slide, image.path, sx(6.05), sy(1.22), sx(6.45), sy(4.9))
        add_caption(slide, image.caption, sx(6.2), sy(6.15), sx(6.1))

    add_slide_number(slide, spec.number, spec.total)


def add_results_slide(prs: Presentation, spec: SlideSpec) -> None:
    slide = prepare_slide(prs, 4, spec.title, remove_non_title=True)

    add_body_box(slide, sx(0.55), sy(1.03), sx(4.25), sy(4.1))
    lines = [line for line in spec.text_lines if clean_text(line)]
    if spec.short_line:
        lines.append("")
        lines.append(spec.short_line)
    add_text_block(
        slide,
        lines,
        sx(0.72),
        sy(1.16),
        sx(3.92),
        sy(3.75),
        font_size=15,
        bold_predicate=lambda text: text.endswith(":"),
    )

    if spec.images:
        image = spec.images[0]
        add_fitted_picture(slide, image.path, sx(5.05), sy(1.18), sx(7.6), sy(3.35))
        add_caption(slide, image.caption, sx(5.12), sy(4.58), sx(7.45))

    if spec.table_rows:
        add_table(slide, spec.table_rows, sx(0.72), sy(5.0), sx(11.9), sy(1.5), font_size=10)

    add_slide_number(slide, spec.number, spec.total)


def add_adaptive_slide(prs: Presentation, spec: SlideSpec) -> None:
    slide = prepare_slide(prs, 4, spec.title, remove_non_title=True)

    add_body_box(slide, sx(0.55), sy(1.02), sx(3.95), sy(5.55))
    add_text_block(
        slide,
        [line for line in spec.text_lines if clean_text(line)],
        sx(0.72),
        sy(1.16),
        sx(3.62),
        sy(4.75),
        font_size=14,
        bold_predicate=lambda text: text.endswith(":"),
    )
    if spec.note:
        emphasize_box(slide, spec.note, sx(0.82), sy(5.95), sx(3.4), sy(0.42))

    if len(spec.images) >= 1:
        add_fitted_picture(slide, spec.images[0].path, sx(4.7), sy(1.16), sx(7.85), sy(2.8))
        add_caption(slide, spec.images[0].caption, sx(4.84), sy(3.95), sx(7.55))
    if len(spec.images) >= 2:
        add_fitted_picture(slide, spec.images[1].path, sx(4.7), sy(4.22), sx(4.0), sy(1.5))
        add_caption(slide, spec.images[1].caption, sx(4.76), sy(5.73), sx(3.88))
    if spec.table_rows:
        add_table(slide, spec.table_rows, sx(8.95), sy(4.28), sx(3.55), sy(1.48), font_size=9)

    add_slide_number(slide, spec.number, spec.total)


def add_final_slide(prs: Presentation, spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[13])
    title_ph = get_title_placeholder(slide)
    if title_ph is not None:
        title_ph.text = clean_text(spec.title)

    body_ph = next((ph for ph in slide.placeholders if ph.placeholder_format.type != PP_PLACEHOLDER.TITLE), None)
    if body_ph is not None:
        tf = body_ph.text_frame
        tf.clear()
        tf.word_wrap = True
        for idx, line in enumerate([line for line in spec.text_lines if clean_text(line)]):
            paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            paragraph.text = clean_text(line)
            paragraph.font.name = BODY_FONT
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = BODY_COLOR
            paragraph.space_after = Pt(3)

    if spec.final_line:
        emphasize_box(slide, spec.final_line, 0.9, 4.85, 7.75, 0.45)

    add_slide_number(slide, spec.number, spec.total)


def build_itmo_presentation(source: Path, template: Path, output: Path) -> None:
    slides = parse_source(source)
    prs = Presentation(template)
    clear_all_template_slides(prs)

    for spec in slides:
        if spec.number == 1:
            add_template_title_slide(prs, spec)
        elif spec.number == 3:
            add_goal_slide(prs, spec)
        elif spec.number == 5:
            add_method_slide(prs, spec)
        elif spec.number == 6:
            add_results_slide(prs, spec)
        elif spec.number == 7:
            add_adaptive_slide(prs, spec)
        elif spec.number == 8:
            add_final_slide(prs, spec)
        else:
            add_text_image_slide(prs, spec)

    prs.save(output)
    print(f"Saved: {output}")


def main() -> None:
    parser = argparse.ArgumentParser(description="Build the pre-defense PPTX using the ITMO template.")
    parser.add_argument("--input", type=Path, default=SOURCE_DEFAULT)
    parser.add_argument("--template", type=Path, default=TEMPLATE_DEFAULT)
    parser.add_argument("--output", type=Path, default=OUTPUT_DEFAULT)
    args = parser.parse_args()

    build_itmo_presentation(args.input, args.template, args.output)


if __name__ == "__main__":
    main()

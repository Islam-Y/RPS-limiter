#!/usr/bin/env python3
from __future__ import annotations

import argparse
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import List

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SOURCE_DEFAULT = Path("/Users/islam/IdeaProjects/RPS-limiter/services/Презентация_предзащита_5мин.txt")
OUTPUT_DEFAULT = Path("/Users/islam/IdeaProjects/RPS-limiter/services/Презентация_предзащита_5мин.pptx")


TITLE_FONT = "Georgia"
BODY_FONT = "Arial"
TITLE_COLOR = RGBColor(255, 255, 255)
BODY_COLOR = RGBColor(31, 41, 55)
MUTED_COLOR = RGBColor(93, 104, 116)
NAVY = RGBColor(17, 45, 78)
TEAL = RGBColor(34, 94, 105)
ACCENT = RGBColor(201, 133, 49)
BG = RGBColor(248, 247, 243)
PANEL = RGBColor(255, 255, 255)
LINE = RGBColor(220, 224, 228)
LIGHT_ACCENT = RGBColor(243, 236, 226)


@dataclass
class ImageSpec:
    path: Path
    caption: str = ""


@dataclass
class SlideSpec:
    number: int
    total: int
    name: str
    title: str = ""
    text_lines: List[str] = field(default_factory=list)
    short_line: str = ""
    final_line: str = ""
    visual_text: str = ""
    images: List[ImageSpec] = field(default_factory=list)
    table_sources: List[Path] = field(default_factory=list)
    table_rows: List[List[str]] = field(default_factory=list)
    note: str = ""
    speaker: str = ""


def clean_text(value: str) -> str:
    return value.replace("`", "").strip()


def parse_source(path: Path) -> List[SlideSpec]:
    lines = path.read_text(encoding="utf-8").splitlines()
    slides: List[SlideSpec] = []
    current_block: List[str] = []

    def flush(block: List[str]) -> None:
        if not block:
            return
        slides.append(parse_slide_block(block))

    for line in lines:
        if line.startswith("Слайд "):
            flush(current_block)
            current_block = [line]
        elif current_block:
            current_block.append(line)
    flush(current_block)
    return slides


def parse_slide_block(block: List[str]) -> SlideSpec:
    header = block[0].strip()
    match = re.match(r"^Слайд\s+(\d+)/(\d+)\.\s*(.+)$", header)
    if not match:
        raise ValueError(f"Не удалось разобрать заголовок блока: {header}")

    spec = SlideSpec(number=int(match.group(1)), total=int(match.group(2)), name=match.group(3).strip())

    marker_map = {
        "Заголовок:": "title",
        "Текст на слайде:": "text_lines",
        "Короткая мысль на слайде:": "short_line",
        "Финальная фраза на слайде:": "final_line",
        "Визуал:": "visual_text",
        "Примечание:": "note",
        "Реплика доклада:": "speaker",
        "Подпись для слайда:": "image_caption",
    }

    current_field = None
    last_image: ImageSpec | None = None
    capture_table = False

    for raw in block[1:]:
        line = raw.rstrip()
        stripped = line.strip()

        if not stripped:
            if current_field == "text_lines":
                spec.text_lines.append("")
            continue

        if stripped.startswith("Место вставки"):
            current_field = None
            capture_table = False
            continue

        if stripped in {"[ВСТАВИТЬ ИЗОБРАЖЕНИЕ]", "[ВСТАВИТЬ КОМПАКТНУЮ ТАБЛИЦУ]", "[БЕЗ ИЗОБРАЖЕНИЯ]"}:
            current_field = None
            continue

        if stripped.startswith("Ссылка на изображение:"):
            image_path = Path(stripped.split(":", 1)[1].strip())
            last_image = ImageSpec(path=image_path)
            spec.images.append(last_image)
            current_field = None
            capture_table = False
            continue

        if stripped.startswith("Ссылка на таблицу-источник:"):
            spec.table_sources.append(Path(stripped.split(":", 1)[1].strip()))
            current_field = None
            capture_table = False
            continue

        if stripped == "Рекомендуемая таблица:":
            capture_table = True
            current_field = None
            continue

        if capture_table and stripped.startswith("|"):
            row = [cell.strip() for cell in stripped.strip("|").split("|")]
            if row and not all(set(cell) <= {"-"} for cell in row):
                spec.table_rows.append(row)
            continue
        elif capture_table and not stripped.startswith("|"):
            capture_table = False

        marker = next((m for m in marker_map if stripped.startswith(m)), None)
        if marker:
            current_field = marker_map[marker]
            value = stripped.split(":", 1)[1].strip()
            if value:
                assign_field(spec, current_field, value, last_image)
            continue

        if current_field:
            assign_field(spec, current_field, stripped, last_image)

    return spec


def assign_field(spec: SlideSpec, field_name: str, value: str, last_image: ImageSpec | None) -> None:
    value = clean_text(value)
    if not value:
        return
    if field_name == "title":
        spec.title = f"{spec.title} {value}".strip()
    elif field_name == "text_lines":
        spec.text_lines.append(value)
    elif field_name == "short_line":
        spec.short_line = f"{spec.short_line} {value}".strip()
    elif field_name == "final_line":
        spec.final_line = f"{spec.final_line} {value}".strip()
    elif field_name == "visual_text":
        spec.visual_text = f"{spec.visual_text} {value}".strip()
    elif field_name == "note":
        spec.note = f"{spec.note} {value}".strip()
    elif field_name == "speaker":
        spec.speaker = f"{spec.speaker} {value}".strip()
    elif field_name == "image_caption" and last_image is not None:
        last_image.caption = f"{last_image.caption} {value}".strip()


def set_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_top_band(slide, title: str) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.72))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(0.72), Inches(13.333), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.12), Inches(11.9), Inches(0.4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = clean_text(title)
    p.font.name = TITLE_FONT
    p.font.size = Pt(25)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR


def add_slide_number(slide, number: int, total: int, *, light: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(12.2), Inches(7.03), Inches(0.7), Inches(0.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{number}/{total}"
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = BODY_FONT
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(240, 243, 246) if light else MUTED_COLOR


def add_body_box(slide, left: float, top: float, width: float, height: float) -> object:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = LINE
    return shape


def add_text_block(
    slide,
    lines: List[str],
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    font_size: int = 18,
    color: RGBColor = BODY_COLOR,
    bold_predicate=None,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    first = True
    for raw_line in lines:
        text = clean_text(raw_line)
        if not text:
            continue
        paragraph = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        paragraph.text = text
        paragraph.font.name = BODY_FONT
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.font.bold = bool(bold_predicate and bold_predicate(text))
        paragraph.space_after = Pt(4)
        paragraph.line_spacing = 1.05


def add_fitted_picture(slide, image_path: Path, left: float, top: float, width: float, height: float) -> None:
    frame = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = PANEL
    frame.line.color.rgb = LINE

    with Image.open(image_path) as img:
        img_width, img_height = img.size

    box_aspect = width / height
    img_aspect = img_width / img_height

    if img_aspect >= box_aspect:
        render_w = width - 0.12
        render_h = render_w / img_aspect
    else:
        render_h = height - 0.12
        render_w = render_h * img_aspect

    render_left = left + (width - render_w) / 2
    render_top = top + (height - render_h) / 2
    slide.shapes.add_picture(
        str(image_path),
        Inches(render_left),
        Inches(render_top),
        width=Inches(render_w),
        height=Inches(render_h),
    )


def add_caption(slide, caption: str, left: float, top: float, width: float) -> None:
    if not caption:
        return
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.28))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = clean_text(caption)
    p.font.name = BODY_FONT
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED_COLOR
    p.alignment = PP_ALIGN.CENTER


def add_table(slide, rows: List[List[str]], left: float, top: float, width: float, height: float, font_size: int = 12) -> None:
    if not rows:
        return

    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table

    col_width = width / len(rows[0])
    for col in table.columns:
        col.width = Inches(col_width)

    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = clean_text(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if r == 0 else (LIGHT_ACCENT if r % 2 == 1 else PANEL)

            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.name = BODY_FONT
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = r == 0
            paragraph.font.color.rgb = TITLE_COLOR if r == 0 else BODY_COLOR
            cell.text_frame.word_wrap = True
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def emphasize_box(slide, text: str, left: float, top: float, width: float, height: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_ACCENT
    shape.line.color.rgb = ACCENT
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = clean_text(text)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = BODY_FONT
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = BODY_COLOR


def build_title_slide(prs: Presentation, spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    left_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.05))
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = TEAL
    left_panel.line.fill.background()

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(1.35), Inches(0.12), Inches(1.9))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.05), Inches(1.2), Inches(10.9), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = clean_text(spec.title)
    p.font.name = TITLE_FONT
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    info_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(1.05),
        Inches(4.25),
        Inches(7.6),
        Inches(1.55),
    )
    info_shape.fill.solid()
    info_shape.fill.fore_color.rgb = BG
    info_shape.line.color.rgb = RGBColor(53, 78, 105)
    tf = info_shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(18)
    tf.margin_top = Pt(10)
    for idx, line in enumerate(spec.text_lines):
        if not clean_text(line):
            continue
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.text = clean_text(line)
        paragraph.font.name = BODY_FONT
        paragraph.font.size = Pt(20 if idx == 0 else 17)
        paragraph.font.bold = idx == 0
        paragraph.font.color.rgb = BODY_COLOR
        paragraph.space_after = Pt(3)

    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(9.3),
        Inches(4.55),
        Inches(2.6),
        Inches(0.72),
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = ACCENT
    tag.line.fill.background()
    tf = tag.text_frame
    p = tf.paragraphs[0]
    p.text = "Предзащита, 5 минут"
    p.font.name = BODY_FONT
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_slide_number(slide, spec.number, spec.total, light=True)


def build_text_image_slide(prs: Presentation, spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_band(slide, spec.title)

    add_body_box(slide, 0.55, 1.02, 5.05, 5.7)
    lines = [line for line in spec.text_lines if clean_text(line)]
    if spec.short_line:
        lines.append("")
        lines.append(spec.short_line)
    add_text_block(
        slide,
        lines,
        0.72,
        1.18,
        4.72,
        5.35,
        font_size=18,
        bold_predicate=lambda text: text.endswith(":"),
    )

    if spec.images:
        image = spec.images[0]
        add_fitted_picture(slide, image.path, 5.85, 1.18, 6.92, 4.95)
        add_caption(slide, image.caption, 5.95, 6.15, 6.72)

    if spec.note:
        emphasize_box(slide, spec.note, 0.82, 6.18, 4.5, 0.48)

    add_slide_number(slide, spec.number, spec.total)


def build_goal_slide(prs: Presentation, spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_band(slide, spec.title)

    goal_lines = []
    tasks = []
    target = goal_lines
    for line in spec.text_lines:
        text = clean_text(line)
        if not text:
            continue
        if text == "Задачи:":
            target = tasks
            continue
        if text == "Цель работы:":
            continue
        target.append(text)

    goal_text = " ".join(goal_lines[:1]) if goal_lines else ""
    emphasize_box(slide, f"Цель работы: {goal_text}", 0.68, 1.15, 11.95, 0.9)

    add_body_box(slide, 0.68, 2.3, 5.85, 3.95)
    add_body_box(slide, 6.8, 2.3, 5.85, 3.95)

    left_tasks = tasks[:4]
    right_tasks = tasks[4:]
    add_text_block(slide, left_tasks, 0.9, 2.55, 5.35, 3.55, font_size=16)
    add_text_block(slide, right_tasks, 7.02, 2.55, 5.35, 3.55, font_size=16)

    hint = "Все задачи связаны с практической проверкой решения на реальном проектном стенде."
    emphasize_box(slide, hint, 2.1, 6.45, 9.0, 0.46)
    add_slide_number(slide, spec.number, spec.total)


def build_method_slide(prs: Presentation, spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_band(slide, spec.title)

    add_body_box(slide, 0.55, 1.05, 5.3, 5.7)
    add_text_block(
        slide,
        [line for line in spec.text_lines if clean_text(line)],
        0.72,
        1.2,
        4.96,
        5.2,
        font_size=17,
        bold_predicate=lambda text: text.endswith(":"),
    )
    emphasize_box(slide, "Норма → атака → восстановление", 0.92, 6.1, 4.55, 0.42)

    if spec.images:
        image = spec.images[0]
        add_fitted_picture(slide, image.path, 6.05, 1.22, 6.45, 4.9)
        add_caption(slide, image.caption, 6.2, 6.15, 6.1)

    add_slide_number(slide, spec.number, spec.total)


def build_results_slide(prs: Presentation, spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_band(slide, spec.title)

    add_body_box(slide, 0.55, 1.03, 4.25, 4.1)
    add_text_block(
        slide,
        [line for line in spec.text_lines if clean_text(line)] + (["", spec.short_line] if spec.short_line else []),
        0.72,
        1.18,
        3.9,
        3.8,
        font_size=17,
        bold_predicate=lambda text: text.endswith(":"),
    )

    if spec.images:
        image = spec.images[0]
        add_fitted_picture(slide, image.path, 5.05, 1.18, 7.6, 3.4)
        add_caption(slide, image.caption, 5.15, 4.64, 7.35)

    if spec.table_rows:
        add_table(slide, spec.table_rows, 0.72, 5.1, 11.9, 1.45, font_size=12)

    add_slide_number(slide, spec.number, spec.total)


def build_adaptive_slide(prs: Presentation, spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_band(slide, spec.title)

    add_body_box(slide, 0.55, 1.02, 3.95, 5.65)
    add_text_block(
        slide,
        [line for line in spec.text_lines if clean_text(line)],
        0.72,
        1.16,
        3.62,
        4.9,
        font_size=16,
        bold_predicate=lambda text: text.endswith(":"),
    )
    if spec.note:
        emphasize_box(slide, spec.note, 0.84, 6.02, 3.35, 0.48)

    if len(spec.images) >= 1:
        add_fitted_picture(slide, spec.images[0].path, 4.7, 1.16, 7.85, 2.8)
        add_caption(slide, spec.images[0].caption, 4.85, 3.98, 7.5)
    if len(spec.images) >= 2:
        add_fitted_picture(slide, spec.images[1].path, 4.7, 4.25, 4.0, 1.55)
        add_caption(slide, spec.images[1].caption, 4.75, 5.82, 3.9)
    if spec.table_rows:
        add_table(slide, spec.table_rows, 8.95, 4.3, 3.55, 1.52, font_size=11)

    add_slide_number(slide, spec.number, spec.total)


def build_conclusion_slide(prs: Presentation, spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_band(slide, spec.title)

    add_body_box(slide, 0.72, 1.2, 12.0, 4.9)
    add_text_block(slide, [line for line in spec.text_lines if clean_text(line)], 0.95, 1.42, 11.5, 4.3, font_size=18)
    if spec.final_line:
        emphasize_box(slide, spec.final_line, 1.0, 6.2, 11.2, 0.5)

    add_slide_number(slide, spec.number, spec.total)


def build_presentation(slides: List[SlideSpec], output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for spec in slides:
        if spec.number == 1:
            build_title_slide(prs, spec)
        elif spec.number == 3:
            build_goal_slide(prs, spec)
        elif spec.number == 5:
            build_method_slide(prs, spec)
        elif spec.number == 6:
            build_results_slide(prs, spec)
        elif spec.number == 7:
            build_adaptive_slide(prs, spec)
        elif spec.number == 8:
            build_conclusion_slide(prs, spec)
        else:
            build_text_image_slide(prs, spec)

    prs.save(output_path)


def main() -> None:
    parser = argparse.ArgumentParser(description="Build PPTX for the 5-minute pre-defense presentation.")
    parser.add_argument("--input", type=Path, default=SOURCE_DEFAULT)
    parser.add_argument("--output", type=Path, default=OUTPUT_DEFAULT)
    args = parser.parse_args()

    slides = parse_source(args.input)
    build_presentation(slides, args.output)
    print(f"Saved: {args.output}")


if __name__ == "__main__":
    main()
